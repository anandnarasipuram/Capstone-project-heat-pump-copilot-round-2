#!/usr/bin/env python3
"""Builds presentation.pptx (the Round 2 required deliverable) from the
structured slide content below. Re-run after editing this file to
regenerate the deck — the output is fully editable in PowerPoint/Keynote
afterwards, this script just gets a consistent-looking first draft.

Usage:
    python3 presentation/build_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation.pptx"

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
DARK_TEAL = RGBColor(0x13, 0x3A, 0x38)
TEAL = RGBColor(0x1F, 0x5C, 0x58)
AMBER = RGBColor(0xE8, 0x9B, 0x2E)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x23, 0x2B, 0x2A)
MUTED = RGBColor(0x5B, 0x66, 0x64)
RED = RGBColor(0xB4, 0x3B, 0x2E)
GREEN = RGBColor(0x3E, 0x7A, 0x52)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color=LIGHT_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=15, color=INK, bullet_color=AMBER, gap=6, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"›  {item}"
        p.space_after = Pt(gap)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
    return box


def header(slide, kicker, title, page_no):
    fill_bg(slide)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), AMBER)
    add_text(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.4), kicker,
              size=13, color=AMBER, bold=True)
    add_text(slide, Inches(0.7), Inches(0.68), Inches(11.5), Inches(0.9), title,
              size=30, color=DARK_TEAL, bold=True)
    add_rect(slide, Inches(0.7), Inches(1.35), Inches(2.2), Pt(3), AMBER)
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35), str(page_no),
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.7), Inches(7.05), Inches(4), Inches(0.35),
              "Heat Pump Copilot · Round 2", size=11, color=MUTED)


def table_slide(kicker, title, page_no, headers, rows, col_widths=None, font_size=13):
    slide = add_slide()
    header(slide, kicker, title, page_no)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    left, top = Inches(0.7), Inches(1.65)
    width, height = Inches(11.9), Inches(0.5 * n_rows)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for c, htext in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_TEAL
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = INK
    return slide


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
slide = add_slide()
fill_bg(slide, DARK_TEAL)
add_rect(slide, 0, Inches(6.9), SLIDE_W, Inches(0.6), AMBER)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5), "HEAT PUMP COPILOT",
          size=20, color=AMBER, bold=True)
add_text(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.6),
          "Field Commissioning & HEMS Connectivity Copilot",
          size=40, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(4.35), Inches(11), Inches(0.6),
          "Round 2 — Consulting Package, MVP & Deployment Plan", size=18, color=WHITE)
add_text(slide, Inches(1), Inches(6.3), Inches(8), Inches(0.5),
          "Anand Narasipuram · Ironhack AI Capstone", size=14, color=RGBColor(0xCF, 0xDD, 0xDA))

# ---------------------------------------------------------------------------
# Slide 2 — Agenda
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "OVERVIEW", "Agenda", 2)
add_bullets(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8), [
    "The problem — installer capacity, not demand",
    "Proposed solution — a three-mode advisory copilot",
    "From POC (Round 1) to a working MVP (Round 2) — live demo",
    "Business case — ROI, break-even, risk matrix",
    "Compliance — EU AI Act classification, GDPR posture",
    "Strategic plan — POC → Pilot → Full Deployment, go-to-market",
    "The ask — what we need to greenlight the pilot",
], size=19, gap=14)

# ---------------------------------------------------------------------------
# Slide 3 — The problem
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "THE PROBLEM", "Two field problems that look identical", 3)
add_bullets(slide, Inches(0.7), Inches(1.9), Inches(6.1), Inches(4.6), [
    "A genuine hardware fault requiring a certified technician",
    "A HEMS connectivity/pairing/app issue that needs no hardware visit at all",
    "Installers cannot reliably tell them apart from symptoms alone — a documented, industry-wide ambiguity, not a Chleo-specific bug",
    "The costliest failure: a false hardware-fault dispatch — an unneeded parts/technician visit",
], size=16, gap=14)
add_rect(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(4.6), WHITE)
add_text(slide, Inches(7.5), Inches(2.1), Inches(4.8), Inches(0.5), "Why the company can't just hire more installers",
          size=15, color=DARK_TEAL, bold=True)
add_bullets(slide, Inches(7.5), Inches(2.7), Inches(4.8), Inches(3.6), [
    "Germany trains ~12,000 SHK tradespeople/yr against ~35,000/yr needed",
    "6M-heat-pump-by-2030 target needs a ~50% installer-workforce uplift",
    "2024 sales: ~193,000 (down 46% vs. target of 500,000/yr)",
    "A capacity multiplier for existing installers beats a hiring plan the labor market can't supply",
], size=14, color=INK, gap=12)

# ---------------------------------------------------------------------------
# Slide 4 — Solution overview
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "THE SOLUTION", "One copilot, three integrated modes", 4)
cards = [
    ("🩺", "Fault Triage Copilot", "Reactive — the flagship", "Fault code / symptom → hardware, connectivity, or installer-error classification → fix guidance or escalation"),
    ("✅", "Commissioning Checker", "Preventive", "Confirms commissioning steps were completed before an installer signs a job off"),
    ("📉", "COP-Drop Early-Warning", "Predictive", "Flags a unit likely to fail before anyone reports a fault, using seasonal COP baselines"),
]
x = Inches(0.7)
for emoji, title, tag, desc in cards:
    add_rect(slide, x, Inches(1.9), Inches(3.85), Inches(4.6), WHITE)
    add_rect(slide, x, Inches(1.9), Inches(3.85), Inches(0.12), AMBER)
    add_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.7), emoji, size=32)
    add_text(slide, x + Inches(0.3), Inches(2.95), Inches(3.3), Inches(0.7), title, size=17, bold=True, color=DARK_TEAL)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(3.3), Inches(0.4), tag, size=12, color=AMBER, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.0), Inches(3.3), Inches(2.3), desc, size=13, color=MUTED)
    x += Inches(4.1)

# ---------------------------------------------------------------------------
# Slide 5 — POC recap
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "ROUND 1 RECAP", "The POC: n8n + Telegram, keyword-grounded", 5)
add_bullets(slide, Inches(0.7), Inches(1.9), Inches(11.6), Inches(4.6), [
    "13 fault codes resolved instantly by deterministic lookup — no API cost, fully auditable",
    "Free-text symptoms grounded in a manual knowledge base via keyword matching (Vaillant + Octopus documentation)",
    "OpenAI classification for anything the lookup table doesn't cover — replies in the installer's own language",
    "Every response labeled: \"AI-suggested triage — confirm before acting\" — advisory only, by design",
    "Live-demoable on Telegram in real time — proved the interaction pattern before investing in a full app",
], size=17, gap=16)

# ---------------------------------------------------------------------------
# Slide 6 — MVP architecture / demo
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "ROUND 2 MVP", "The upgrade: real RAG, one working app", 6)
add_text(slide, Inches(0.7), Inches(1.85), Inches(5.7), Inches(0.4), "What changed", size=15, bold=True, color=DARK_TEAL)
add_bullets(slide, Inches(0.7), Inches(2.3), Inches(5.7), Inches(4.2), [
    "Keyword match → OpenAI embeddings + Pinecone vector search (real RAG)",
    "Telegram-only → a single Streamlit app, all 3 modes, one command to run",
    "LangSmith placeholder → every interaction traced live, not a one-off script",
    "16 offline unit tests — deterministic logic verified with zero API keys",
    "Fails soft everywhere: missing keys degrade to a labeled fallback, never a crash",
], size=14, gap=12)
add_rect(slide, Inches(6.8), Inches(1.85), Inches(5.8), Inches(4.7), WHITE)
add_text(slide, Inches(7.1), Inches(2.05), Inches(5.2), Inches(0.4), "Pipeline (Mode 1)", size=15, bold=True, color=DARK_TEAL)
pipeline = [
    "Installer message",
    "↓  fault-code regex lookup (instant, free)",
    "↓  no match → Pinecone similarity search",
    "     over embedded manual knowledge base",
    "↓  OpenAI classification, grounded in the",
    "     retrieved excerpt, JSON-mode output",
    "↓  structured result + \"confirm before acting\"",
]
add_text(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(3.7), "\n".join(pipeline),
          size=14, color=INK, line_spacing=1.3)

# ---------------------------------------------------------------------------
# Slide 7 — Live demo
# ---------------------------------------------------------------------------
slide = add_slide()
fill_bg(slide, TEAL)
add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(1), "Live Demo",
          size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
          "mvp/app.py — Fault Triage Copilot · Commissioning Checker · COP-Drop Early-Warning",
          size=16, color=RGBColor(0xCF, 0xDD, 0xDA), align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 8 — Success criteria
# ---------------------------------------------------------------------------
slide = table_slide("BUSINESS CASE", "Success criteria (pilot targets)", 8,
    ["Metric", "Baseline (Round 1 synthetic data)", "Pilot target"],
    [
        ["First-visit fix rate", "69.1%", "≥ 79.1% (+10pp)"],
        ["False hardware-fault rate", "10.9%", "≤ 5.9% (−5pp)"],
        ["Median time-to-classification", "n/a (manual triage)", "< 30 seconds"],
    ], col_widths=[Inches(4.5), Inches(4), Inches(3.4)])

# ---------------------------------------------------------------------------
# Slide 9 — ROI
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "BUSINESS CASE", "ROI — central case", 9)
metrics = [
    ("Year 1 ROI", "≈ 241%"),
    ("36-month ROI", "≈ 1,198%"),
    ("Break-even", "~3.5–4 months into pilot"),
    ("Combined build cost", "≈ €20,150"),
]
x = Inches(0.7)
for label, val in metrics:
    add_rect(slide, x, Inches(1.9), Inches(2.85), Inches(2.0), WHITE)
    add_text(slide, x + Inches(0.2), Inches(2.1), Inches(2.45), Inches(0.9), val, size=26, bold=True, color=DARK_TEAL)
    add_text(slide, x + Inches(0.2), Inches(3.15), Inches(2.45), Inches(0.6), label, size=13, color=MUTED)
    x += Inches(3.0)
add_text(slide, Inches(0.7), Inches(4.3), Inches(11.6), Inches(0.5),
          "ROI = (Net Benefit / Total Cost) × 100 — full assumptions table in roi_risk_assessment.md",
          size=13, color=MUTED, italic=True)
add_bullets(slide, Inches(0.7), Inches(5.0), Inches(11.6), Inches(1.8), [
    "Value = avoided false-hardware-fault dispatches + avoided second visits, at a stated 70% realization factor",
    "Figures are gross operational cost avoidance from a synthetic-data baseline — pilot-validation targets, not audited financials",
], size=14, color=INK, gap=10)

# ---------------------------------------------------------------------------
# Slide 10 — Risk matrix
# ---------------------------------------------------------------------------
slide = table_slide("BUSINESS CASE", "Top risks (full matrix: 8 risks, 4 categories)", 10,
    ["Risk", "Category", "L", "I", "Mitigation"],
    [
        ["EU AI Act boundary risk", "Regulatory", "2", "5", "Advisory-only positioning enforced by design + legal gate"],
        ["LLM misclassification", "Technical", "3", "4", "Deterministic-first routing, safe-default escalation, human confirmation"],
        ["Automation bias / over-reliance", "Ethical", "3", "4", "Persistent disclaimer, training, override-rate KPI"],
        ["Low installer adoption", "Operational", "3", "4", "Champions, adoption KPI, low-friction chat interaction"],
    ], col_widths=[Inches(3.3), Inches(1.7), Inches(0.5), Inches(0.5), Inches(5.9)], font_size=12)

# ---------------------------------------------------------------------------
# Slide 11 — EU AI Act
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "COMPLIANCE", "EU AI Act — Limited Risk", 11)
add_rect(slide, Inches(0.7), Inches(1.9), Inches(4.4), Inches(1.0), GREEN)
add_text(slide, Inches(0.9), Inches(2.1), Inches(4.0), Inches(0.7), "LIMITED RISK", size=22, bold=True, color=WHITE)
add_bullets(slide, Inches(0.7), Inches(3.1), Inches(5.8), Inches(3.4), [
    "Not prohibited (Art. 5) — no manipulation, no biometric data",
    "Not high-risk under Annex I — not a safety component of the appliance",
    "Not high-risk under Annex III — reasoned explicitly against the closest calls (critical infrastructure, employment)",
    "→ Only obligation: Art. 50 transparency disclosure — already implemented",
], size=14, gap=12)
add_rect(slide, Inches(6.9), Inches(1.9), Inches(5.7), Inches(4.6), WHITE)
add_text(slide, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.5), "The boundary that would change this", size=15, bold=True, color=DARK_TEAL)
add_text(slide, Inches(7.2), Inches(2.7), Inches(5.1), Inches(3.5),
          "If a future version were wired to auto-trigger a shutdown/lockout, bypass a technician's confirmation, "
          "or evaluate individual installer performance — reclassification is re-run before any such change, not assumed to still hold.",
          size=14, color=INK, line_spacing=1.3)

# ---------------------------------------------------------------------------
# Slide 12 — GDPR
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "COMPLIANCE", "GDPR — what changes at pilot", 12)
add_bullets(slide, Inches(0.7), Inches(1.9), Inches(11.6), Inches(4.6), [
    "Round 1 + Round 2 MVP: public/synthetic data only — no real personal data processed today",
    "Pilot introduces installer-identifiable data (chat content, job records) — legal basis: legitimate interest",
    "Short DPIA completed on the highest-risk activity (free-text chat) — residual risk: low, DPO review recommended before go-live",
    "Art. 22 (automated decision-making) does not apply — the system classifies equipment, not people",
    "Third-party transfers: OpenAI, Pinecone, LangSmith — all have EU-region or SCC-based options; config, not redesign",
], size=16, gap=16)

# ---------------------------------------------------------------------------
# Slide 13 — Strategic plan / timeline
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "STRATEGIC PLAN", "POC → Pilot → Full Deployment", 13)
phases = [
    ("Phase 0", "POC", "Round 1 — done"),
    ("Phase 1", "MVP / Internal\nValidation", "Round 2 — done"),
    ("Phase 2", "Pilot", "Oct 2026 – Jan 2027\n12 weeks, real installers"),
    ("Phase 3", "Full\nDeployment", "Mar – Aug 2027\nphased by region"),
    ("Phase 4", "Scale", "2028+\noptional"),
]
x = Inches(0.5)
w = Inches(2.42)
for i, (num, title, sub) in enumerate(phases):
    color = TEAL if i < 2 else (AMBER if i == 2 else MUTED)
    add_rect(slide, x, Inches(2.2), w, Inches(0.5), color)
    add_text(slide, x, Inches(2.25), w, Inches(0.4), num, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(2.75), w, Inches(2.7), WHITE)
    add_text(slide, x + Inches(0.15), Inches(2.95), w - Inches(0.3), Inches(1.0), title, size=15, bold=True, color=DARK_TEAL)
    add_text(slide, x + Inches(0.15), Inches(4.1), w - Inches(0.3), Inches(1.2), sub, size=11.5, color=MUTED)
    if i < len(phases) - 1:
        add_text(slide, x + w - Inches(0.05), Inches(2.2), Inches(0.5), Inches(0.5), "→", size=20, bold=True, color=AMBER)
    x += w + Inches(0.08)
add_text(slide, Inches(0.7), Inches(5.8), Inches(11.6), Inches(1.0),
          "Pilot → Full Deployment greenlight requires ALL of: ≥7pp first-visit-fix gain, ≥3pp false-hardware-fault "
          "reduction, ≥60% sustained adoption, zero compliance incidents, net-positive installer feedback.",
          size=13, color=INK, italic=True)

# ---------------------------------------------------------------------------
# Slide 14 — GTM / commercialisation
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "GO-TO-MARKET", "From internal tool to (optional) product", 14)
add_text(slide, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.4), "Phases 1–3: internal cost avoidance", size=15, bold=True, color=DARK_TEAL)
add_bullets(slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(3.8), [
    "Buyer: Chleo, funding the field-ops budget",
    "Users: own + partner SHK installers — no external sale needed",
    "Differentiator: grounded in Chleo's own manuals, deterministic-first for speed and cost, advisory-only for a light compliance footprint",
], size=14, gap=12)
add_text(slide, Inches(6.9), Inches(1.85), Inches(5.6), Inches(0.4), "Phase 4: optional licensing", size=15, bold=True, color=DARK_TEAL)
add_bullets(slide, Inches(6.9), Inches(2.3), Inches(5.6), Inches(3.8), [
    "The installer-shortage problem is EU-wide, not Germany-only",
    "White-label to other small-medium manufacturers facing the same constraint",
    "Flat onboarding fee + €15–25/installer/month — priced well below the per-installer value shown in the ROI model",
], size=14, gap=12)

# ---------------------------------------------------------------------------
# Slide 15 — The ask
# ---------------------------------------------------------------------------
slide = add_slide()
header(slide, "THE ASK", "What we need to greenlight the pilot", 15)
add_bullets(slide, Inches(0.7), Inches(1.9), Inches(11.6), Inches(4.6), [
    "Sponsor sign-off to proceed from internal validation to a 12-week pilot",
    "DPO/legal review of the short DPIA before pilot go-live",
    "10–15 pilot installers recruited — a deliberate mix of own and partner SHK",
    "Budget: ≈ €20,150 combined build cost (already spent through this package) + ≈ €150–300/month pilot run cost",
    "A decision-gate meeting at week 12 against the KPIs in strategic_plan.md — not an open-ended commitment",
], size=17, gap=16)

# ---------------------------------------------------------------------------
# Slide 16 — Thank you
# ---------------------------------------------------------------------------
slide = add_slide()
fill_bg(slide, DARK_TEAL)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), AMBER)
add_text(slide, Inches(1), Inches(3.1), Inches(11), Inches(1), "Thank you",
          size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6), "Questions & discussion",
          size=18, color=RGBColor(0xCF, 0xDD, 0xDA), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
          "github.com/anandnarasipuram/Capstone-project-heat-pump-copilot-round-2",
          size=12, color=RGBColor(0xCF, 0xDD, 0xDA), align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")
